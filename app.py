import streamlit as st
import matplotlib.pyplot as plt
import os
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime

# Create output folder if it doesn't exist
if not os.path.exists("output_slides"):
    os.makedirs("output_slides")

st.set_page_config(page_title="Budget Pillars Performance Review", layout="centered")

# Page styling
st.markdown("""
    <style>
        body {
            background-color: #0b1e3f;
            color: white;
        }
        .stApp {
            background-color: #0b1e3f;
        }
        h1, h2, h3, h4, h5, h6 {
            color: white;
        }
        .streamlit-expanderHeader {
            color: white;
        }
    </style>
""", unsafe_allow_html=True)

st.title("Assessing Y2024 Performance and Charting the Course for Y2025")

# Pillar input
pillar = st.selectbox("Select Pillar", [
    "Effective Governance",
    "Human Centric City",
    "Modern Infrastructure",
    "Thriving Economy"
])

st.markdown("### Project Status Input")
total_projects = st.number_input("Total number of projects tracked", min_value=0)
not_commenced = st.number_input("Projects yet to commence", min_value=0, max_value=total_projects)
initiation = st.number_input("Projects at initiation", min_value=0, max_value=total_projects)
in_progress = st.number_input("Projects in progress", min_value=0, max_value=total_projects)
completed = st.number_input("Projects completed", min_value=0, max_value=total_projects)

not_completed = total_projects - completed

st.markdown("### Progress Ratings")
green = st.number_input("Projects between 80% - 100% (Green)", min_value=0, max_value=total_projects)
amber = st.number_input("Projects between 60% - 79% (Amber)", min_value=0, max_value=total_projects)
red = st.number_input("Projects between 0% - 59% (Red)", min_value=0, max_value=total_projects)

if st.button("Generate Analysis") and total_projects > 0:
    # Chart
    fig, ax = plt.subplots()
    ax.pie([green, amber, red], labels=["Green", "Amber", "Red"], autopct='%1.1f%%', colors=['#2ecc71', '#f1c40f', '#e74c3c'])
    ax.set_title("Project Performance Breakdown", color='white')
    fig.patch.set_facecolor('#0b1e3f')
    st.pyplot(fig)

    implications = []
    outlook_recommendations = []

    # Improved implication logic
    green_ratio = green / total_projects
    amber_ratio = amber / total_projects
    red_ratio = red / total_projects

    if green_ratio >= 0.5:
        implications.append("A significant portion of projects made good progress with minimal delays, indicating effective execution in key areas.")
    elif green_ratio < 0.3:
        implications.append("Few projects reached full or near completion, suggesting a need for stronger project management and delivery push.")

    if amber_ratio > 0.3:
        implications.append("Many projects experienced slight delays, which may impact timely achievement of the pillar if not addressed.")

    if red_ratio >= 0.2:
        implications.append("A concerning number of projects are severely delayed and require urgent intervention to recover delivery timelines.")

    # Recommendations based on pillars
    if pillar == "Effective Governance":
        outlook_recommendations = [
            "Fast-track digitization and e-governance tools to support efficiency.",
            "Deploy performance dashboards for continuous MDA tracking.",
            "Institutionalize citizen feedback mechanisms."
        ]
    elif pillar == "Human Centric City":
        outlook_recommendations = [
            "Scale programmes with direct impact on health, education, and social welfare.",
            "Enhance inter-agency collaboration for holistic human development.",
            "Focus on inclusive delivery for vulnerable populations."
        ]
    elif pillar == "Modern Infrastructure":
        outlook_recommendations = [
            "Accelerate delivery timelines for mobility and housing projects.",
            "Improve contractor supervision frameworks.",
            "Enforce regular progress audits for major infrastructure."
        ]
    elif pillar == "Thriving Economy":
        outlook_recommendations = [
            "Strengthen SME support schemes and innovation hubs.",
            "Enable productive infrastructure to boost competitiveness.",
            "Align job creation efforts with local economic strengths."
        ]
    
    # Precompute join strings to avoid using backslashes in f-string expressions.
    implications_text = "\n    ".join(implications)
    outlook_text = "\n    ".join(outlook_recommendations)
    
    expert_text = f"""
    ### Expert Analysis

    **Pillar:** {pillar}  
    **Total Projects Tracked:** {total_projects}  
    - Completed: {completed} ({(completed/total_projects)*100:.1f}%)  
    - Not Completed: {not_completed} ({(not_completed/total_projects)*100:.1f}%)  

    **Performance Breakdown:**  
    - Green: {green}  
    - Amber: {amber}  
    - Red: {red}  

    **Implications:**  
    {implications_text}

    **2025 Outlook & Recommendations:**  
    {outlook_text}
    """

    st.markdown(expert_text)

    # Save as PNG (Summary with all text white)
    slide = Image.new("RGB", (1000, 800), color="#0b1e3f")
    draw = ImageDraw.Draw(slide)
    font = ImageFont.load_default()

    slide_text = f"{pillar}\n\nProjects: {total_projects}\nCompleted: {completed}\nNot Completed: {not_completed}\n\nGreen: {green}\nAmber: {amber}\nRed: {red}\n\nImplications:\n- " + "\n- ".join(implications) + "\n\n2025 Outlook:\n- " + "\n- ".join(outlook_recommendations)
    draw.text((30, 30), slide_text, fill="white", font=font)

    png_filename = f"output_slides/{pillar.replace(' ', '_')}_summary.png"
    slide.save(png_filename)
    st.success("PNG slide summary saved!")
    with open(png_filename, "rb") as f:
        st.download_button("Download Slide Summary (PNG)", f, file_name=os.path.basename(png_filename), mime="image/png")

    # Save as PPTX (Make sure text is white)
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]
    ppt_slide = prs.slides.add_slide(slide_layout)
    title_box = ppt_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = f"{pillar} – Performance Summary"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)  # Ensure title text is white

    content = f"Projects Tracked: {total_projects}\nCompleted: {completed}\nNot Completed: {not_completed}\n\nGreen: {green}\nAmber: {amber}\nRed: {red}\n\nImplications:\n- " + "\n- ".join(implications) + "\n\n2025 Outlook:\n- " + "\n- ".join(outlook_recommendations)

    content_box = ppt_slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.text = content
    for para in tf.paragraphs:
        para.font.size = Pt(16)
        para.font.color.rgb = RGBColor(255, 255, 255)  # Ensure content text is white

    pptx_filename = f"output_slides/{pillar.replace(' ', '_')}_summary.pptx"
    prs.save(pptx_filename)
    st.success("PPTX slide summary saved!")
    with open(pptx_filename, "rb") as f:
        st.download_button("Download Slide Summary (PPTX)", f, file_name=os.path.basename(pptx_filename), mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
